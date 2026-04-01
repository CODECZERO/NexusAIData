"""
Nexus AI v4.0 — Executive Report Service
Synthesis of cross-tab analysis into conversational executive summaries and PPTX exports.
"""

from __future__ import annotations
import io
import os
from datetime import datetime
from typing import Any, Dict, List, Optional
from loguru import logger

from models import (
    FullAnalysisResult, 
    ExecutiveSummary, 
    ExecutiveSummarySection,
    QualitySeverity
)

class ExecutiveReportService:
    """Service for generating high-level executive insights and professional presentation decks."""

    async def generate_executive_summary(self, analysis: FullAnalysisResult) -> ExecutiveSummary:
        """Synthesize a meta-narrative from across all analysis domains."""
        logger.info(f"Synthesizing executive summary for session: {analysis.session_id}")
        
        # 1. Key Takeaway (Meta-insight)
        takeaway = "Overall, the dataset shows strong potential for optimization, though data quality issues in key dimensions should be addressed before final deployment."
        if analysis.quality and analysis.quality.overall_score < 70:
            takeaway = "CRITICAL: Significant quality issues detected. High risk of biased insights. Immediate cleansing required."
        elif analysis.ml_results and analysis.ml_results.forecast and analysis.ml_results.forecast.best_mape < 10:
            takeaway = "The business shows high predictability with clear seasonal trends. Expansion is recommended based on forecast stability."

        # 2. Strategic Pillars
        pillars = []
        
        # Data Health Pillar
        health_status = "Good" if analysis.quality and analysis.quality.overall_score > 80 else "Variable"
        pillars.append(ExecutiveSummarySection(
            title="Operational Data Integrity",
            content=f"Data health is currently {health_status}. Detected {analysis.quality.total_issues if analysis.quality else 0} total quality issues across key functional areas.",
            impact_metrics={
                "Quality Score": f"{analysis.quality.overall_score if analysis.quality else 0}%",
                "Critical Issues": str(len(analysis.quality.critical_issues) if analysis.quality else 0)
            }
        ))

        # Predictive Pillar
        forecast_summary = "Forecast unavailable."
        if analysis.ml_results and analysis.ml_results.forecast:
            f = analysis.ml_results.forecast
            forecast_summary = f"Detected significant {f.best_model} patterns. Predictability is high (MAPE: {f.best_mape if f.best_mape else 'N/A'}%)."
        
        pillars.append(ExecutiveSummarySection(
            title="Growth & Predictability",
            content=forecast_summary,
            impact_metrics={
                "Forecast Accuracy": f"{(100 - f.best_mape) if analysis.ml_results and analysis.ml_results.forecast and f.best_mape else 0}%",
                "Primary Driver": analysis.ml_results.feature_importance.top_5_drivers[0][0] if analysis.ml_results and analysis.ml_results.feature_importance and analysis.ml_results.feature_importance.top_5_drivers else "Unknown"
            }
        ))

        # Efficiency Pillar (Simulation)
        sim_summary = "Simulation model baseline established."
        if analysis.simulation:
            s = analysis.simulation
            sim_summary = f"Identified {len(s.constrained_results)} optimal paths. Highest predicted improvement is {max([r.improvement_pct for r in s.constrained_results]) if s.constrained_results else 0}%."

        pillars.append(ExecutiveSummarySection(
            title="Strategic Optimization",
            content=sim_summary,
            impact_metrics={
                "Optimized ROI": f"{max([r.predicted_outcome for r in s.constrained_results]) if analysis.simulation and s.constrained_results else 0:,.2f}",
                "Scenario Count": str(len(analysis.simulation.scenarios) if analysis.simulation else 0)
            }
        ))

        # 3. Risk Assessment
        risk = "Moderate environmental variance."
        if analysis.quality and len(analysis.quality.critical_issues) > 0:
            risk = "HIGH: Data gaps in critical levers may lead to suboptimal decision-making."
        
        # 4. Next Steps
        steps = [
            "Validate top forecast drivers with domain experts.",
            "Deploy automated cleaning pipelines for high-cardinality columns.",
            "Schedule follow-up simulation for Q3 budget planning."
        ]

        return ExecutiveSummary(
            session_id=analysis.session_id,
            key_takeaway=takeaway,
            strategic_pillars=pillars,
            risk_assessment=risk,
            next_steps=steps
        )

    async def export_to_pptx(self, analysis: FullAnalysisResult, summary: ExecutiveSummary) -> bytes:
        """Generate a professional PPTX deck. Requires python-pptx."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            logger.warning("python-pptx not installed. Returning empty buffer.")
            return b""

        prs = Presentation()
        
        # ── Slide 1: Title ──────────────────
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"Executive Intelligence Report: {analysis.filename}"
        subtitle.text = f"Generated by Nexus AI • {datetime.now().strftime('%B %d, %Y')}"

        # ── Slide 2: Executive Summary ───────
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Executive Summary"
        
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = summary.key_takeaway
        
        for pillar in summary.strategic_pillars:
            p = tf.add_paragraph()
            p.text = f"{pillar.title}: {pillar.content}"
            p.level = 1

        # ── Slide 3: Operational Health ──────
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = "Operational Data Health"
        
        if analysis.quality:
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            
            table = slide.shapes.add_table(len(analysis.quality.critical_issues) + 1, 3, left, top, width, height).table
            table.columns[0].width = Inches(2)
            table.columns[1].width = Inches(4)
            table.columns[2].width = Inches(2)
            
            table.cell(0, 0).text = "Severity"
            table.cell(0, 1).text = "Issue"
            table.cell(0, 2).text = "Impact"
            
            for i, issue in enumerate(analysis.quality.critical_issues[:5]):
                table.cell(i+1, 0).text = "CRITICAL"
                table.cell(i+1, 1).text = issue.title
                table.cell(i+1, 2).text = f"{issue.rows_affected_pct:.1f}%"

        # ── Slide 4: Forecast Projections ────
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = "Strategic Forecast Projections"
        # In a real implementation, we would insert an image of the Plotly chart here.
        # prs.slides[3].shapes.add_picture('forecast.png', left, top)
        
        # ── Slide 5: Next Steps ──────────────
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = "Strategic Next Steps"
        tf = slide.shapes.placeholders[1].text_frame
        for step in summary.next_steps:
            p = tf.add_paragraph()
            p.text = step
            p.bullet = True

        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()
